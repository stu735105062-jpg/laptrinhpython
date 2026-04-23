from pptx import Presentation
from pptx.util import Inches

# Tạo presentation mới
prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Quản Lý Công Việc Chuyên Nghiệp"
subtitle.text = "Ứng dụng Python với Tkinter và SQLite\nTác giả: [Your Name]\nNgày: 23/04/2026"

# Slide 2: Giới thiệu
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Giới Thiệu Dự Án'
tf = body_shape.text_frame
tf.text = 'Đây là một ứng dụng quản lý công việc được viết bằng Python, sử dụng giao diện đồ họa Tkinter với thư viện ttkbootstrap để tạo giao diện hiện đại. Ứng dụng cho phép người dùng thêm, xem, cập nhật và xóa các công việc với deadline, thời gian, và mức độ quan trọng.'

# Slide 3: Tính Năng Chính
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Tính Năng Chính'
tf = body_shape.text_frame
tf.text = '• Thêm công việc mới: Sử dụng form nhập liệu với tên công việc, mô tả, deadline, thời gian, và đánh dấu quan trọng.\n• Hiển thị danh sách công việc: Lọc theo trạng thái (Tất cả, Hoàn thành, Chưa hoàn thành).\n• Thông báo: Tự động thông báo khi đến giờ làm việc.\n• Lịch: Chọn ngày deadline bằng calendar picker.\n• Cơ sở dữ liệu: Lưu trữ dữ liệu trong SQLite.'

# Slide 4: Cấu Trúc Dự Án
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Cấu Trúc Dự Án'
tf = body_shape.text_frame
tf.text = '• main.py: File chính, chứa giao diện chính và logic ứng dụng.\n• database.py: Xử lý cơ sở dữ liệu SQLite.\n• add_work.py: Form thêm công việc mới.\n• calender.py: Dialog chọn ngày.\n• test.py: File test (có thể là ví dụ đọc file).\n• README.md: Tài liệu mô tả dự án.\n• diagram.puml: Sơ đồ UML PlantUML.'

# Slide 5: Sơ Đồ UML
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Sơ Đồ UML'
tf = body_shape.text_frame
tf.text = 'Sơ đồ lớp UML mô tả các thành phần chính:\n\n@startuml\nclass DataEntryForm {\n    +work_name: StringVar\n    +create_form_entry(label, variable)\n    ...\n}\n\nclass DatePickerDialog {\n    +parent: Misc\n    +_setup_calendar()\n}\n\nclass Database {\n    +connect()\n    +write(task_data)\n    ...\n}\n\nMain -- DataEntryForm : uses\n...\n@enduml\n\n(Xem chi tiết trong file diagram.puml)'

# Slide 6: Cài Đặt và Chạy
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Cài Đặt và Chạy'
tf = body_shape.text_frame
tf.text = 'Yêu cầu: Python 3.x, tkinter, ttkbootstrap, windows_toasts\n\n1. Cài đặt thư viện:\n   pip install ttkbootstrap windows_toasts\n\n2. Chạy ứng dụng:\n   python main.py'

# Slide 7: Kết Luận
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Kết Luận'
tf = body_shape.text_frame
tf.text = 'Ứng dụng quản lý công việc đơn giản nhưng hiệu quả, sử dụng Python và SQLite. Giao diện thân thiện với người dùng nhờ ttkbootstrap. Có thể mở rộng thêm tính năng như đồng bộ đám mây hoặc đa nền tảng.'

# Lưu file
prs.save('presentation.pptx')
print("PowerPoint presentation created: presentation.pptx")